import re
import io
import os
import textwrap
from pptx import Presentation
from pptx.util import Pt

def intelligent_sentence_split(text):
    """
    Intelligently split text into sentences, preserving all punctuation.
    
    Key improvements:
    1. Use advanced regex to split sentences correctly
    2. Preserve all original punctuation
    3. Ensure each sentence starts with a capital letter
    """
    # Regex that splits on sentence-ending punctuation while keeping the punctuation
    sentence_pattern = r'(?<=[.!?])\s+'
    
    # Split sentences, preserving punctuation
    sentences = re.split(sentence_pattern, text)
    
    # Clean and prepare sentences
    cleaned_sentences = []
    for sentence in sentences:
        # Trim whitespace, but keep original punctuation
        sentence = sentence.strip()
        
        # Capitalize first letter if not already capitalized
        if sentence and not sentence[0].isupper():
            sentence = sentence[0].upper() + sentence[1:]
        
        if sentence:
            cleaned_sentences.append(sentence)
    
    return cleaned_sentences

def clean_text(text):
    """
    Enhanced text cleaning that preserves all meaningful punctuation.
    
    Goals:
    1. Remove only truly extraneous characters
    2. Retain all standard punctuation marks
    3. Normalize whitespace without losing context
    """
    if isinstance(text, list):
        text = "\n".join(text)
    
    # Remove specific unwanted characters while keeping essential punctuation
    # This regex keeps letters, numbers, spaces, and standard punctuation
    text = re.sub(r'[*#]', '', text)
    
    # Normalize whitespace, but carefully
    text = re.sub(r'\s+', ' ', text).strip()
    text = re.sub(r'\n+', '\n', text)
    
    return text

def add_slide(prs, title, content):
    """
    Enhanced slide creation with intelligent bullet point generation.
    
    Improvements:
    1. Create bullet points based on complete sentences
    2. Preserve original punctuation
    3. Ensure readability and professional formatting
    """
    slide_layout = prs.slide_layouts[1]  # Title & Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Consistent title formatting
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)

    # Prepare content area
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    # Intelligent bullet point generation
    for paragraph in content:
        # Split into sentences while preserving punctuation
        sentences = intelligent_sentence_split(paragraph)
        
        for sentence in sentences:
            # Add each sentence as a separate bullet point
            p = text_frame.add_paragraph()
            p.text = sentence
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            p.level = 0  # PowerPoint's default bullet style

def create_presentation(file_texts):
    """
    Create a PowerPoint presentation with enhanced text processing.
    
    Key Features:
    1. Preserve original text structure
    2. Intelligent sentence and bullet point generation
    3. Consistent formatting
    """
    # Use existing template or create new presentation
    template_path = "Ion.pptx"
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    # Add Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Vehicle Rental System - Summary Presentation"
    slide.placeholders[1].text = "Created by AI PPT Generator"

    # Processing parameters
    max_lines_per_slide = 6
    max_chars_per_line = 80

    for filename, text in file_texts.items():
        # Clean the text while preserving structure
        cleaned_text = clean_text(text)
        
        # Add a file title slide
        file_title_slide = prs.slides.add_slide(prs.slide_layouts[5])
        file_title_slide.shapes.title.text = f"Summary of {filename}"

        # Prepare slide content
        current_slide_text = []
        
        # Split text into manageable chunks
        wrapped_lines = textwrap.wrap(cleaned_text, width=max_chars_per_line)
        
        for wrapped_line in wrapped_lines:
            if wrapped_line.strip():
                current_slide_text.append(wrapped_line.strip())

                # Add a new slide when current slide is full
                if len(current_slide_text) == max_lines_per_slide:
                    add_slide(prs, f"Key Points from {filename}", current_slide_text)
                    current_slide_text = []

        # Add any remaining text to a final slide
        if current_slide_text:
            add_slide(prs, f"Additional Info from {filename}", current_slide_text)

    # Save presentation
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io
